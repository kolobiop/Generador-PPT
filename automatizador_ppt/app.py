
from flask import Flask, render_template, request, send_file

app = Flask(__name__)



import tkinter as tk
from tkinter import filedialog, messagebox
from pptx import Presentation
from pptx.util import Inches
from datetime import datetime
import os

def generar_ppt(foto_ruta, descripcion, responsable, nombre_archivo):
    # Validar que todos los campos estén llenos
    if not foto_ruta or not descripcion or not responsable or not nombre_archivo:
        messagebox.showerror("Error", "Todos los campos son obligatorios.")
        return

    # Crear una presentación
    prs = Presentation()

    # Agregar una diapositiva con layout vacío
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    # Agregar la imagen en la parte izquierda
    try:
        slide.shapes.add_picture(foto_ruta, Inches(0.5), Inches(1), width=Inches(4))
    except FileNotFoundError:
        messagebox.showerror("Error", "La foto seleccionada no se encuentra.")
        return

    # Agregar la tabla a la derecha de la imagen
    rows = 3
    cols = 2
    left = Inches(5.5)
    top = Inches(1)
    width = Inches(4)
    height = Inches(2)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # Configurar el tamaño de las columnas
    table.columns[0].width = Inches(1.5)
    table.columns[1].width = Inches(2.5)

    # Llenar la tabla con los datos
    table.cell(0, 0).text = "Descripción"
    table.cell(0, 1).text = descripcion

    table.cell(1, 0).text = "Responsable"
    table.cell(1, 1).text = responsable

    # Obtener la fecha y hora actuales
    hora_reparacion = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    table.cell(2, 0).text = "Fecha y Hora"
    table.cell(2, 1).text = hora_reparacion

    # Comprobar si el archivo ya existe
    nombre_archivo_ppt = f"{nombre_archivo}.pptx"
    if os.path.exists(nombre_archivo_ppt):
        if not messagebox.askyesno("Confirmación", f"El archivo '{nombre_archivo_ppt}' ya existe. ¿Deseas sobrescribirlo?"):
            return

    # Guardar la presentación con el nombre personalizado
    prs.save(nombre_archivo_ppt)
    messagebox.showinfo("Éxito", f'PPT generado: {nombre_archivo_ppt}')

def seleccionar_foto():
    file_path = filedialog.askopenfilename(filetypes=[("Image files", "*.jpg;*.jpeg;*.png")])
    entrada_foto.delete(0, tk.END)
    entrada_foto.insert(0, file_path)

# Crear la ventana principal
root = tk.Tk()
root.title("Generador de PPT")

# Etiqueta y entrada para la foto
tk.Label(root, text="Foto de la reparación:").grid(row=0, column=0, padx=10, pady=5)
entrada_foto = tk.Entry(root, width=40)
entrada_foto.grid(row=0, column=1, padx=10, pady=5)
boton_foto = tk.Button(root, text="Seleccionar", command=seleccionar_foto)
boton_foto.grid(row=0, column=2, padx=10, pady=5)

# Etiqueta y entrada para la descripción
tk.Label(root, text="Descripción:").grid(row=1, column=0, padx=10, pady=5)
entrada_descripcion = tk.Entry(root, width=40)
entrada_descripcion.grid(row=1, column=1, padx=10, pady=5)

# Etiqueta y entrada para el responsable
tk.Label(root, text="Responsable:").grid(row=2, column=0, padx=10, pady=5)
entrada_responsable = tk.Entry(root, width=40)
entrada_responsable.grid(row=2, column=1, padx=10, pady=5)

# Etiqueta y entrada para el nombre del archivo PPT
tk.Label(root, text="Nombre del archivo PPT:").grid(row=3, column=0, padx=10, pady=5)
entrada_nombre_archivo = tk.Entry(root, width=40)
entrada_nombre_archivo.grid(row=3, column=1, padx=10, pady=5)

# Botón para generar el PPT
boton_generar = tk.Button(root, text="Generar PPT", command=lambda: generar_ppt(
    entrada_foto.get(),
    entrada_descripcion.get(),
    entrada_responsable.get(),
    entrada_nombre_archivo.get() or f"reporte_reparacion_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
))
boton_generar.grid(row=4, column=1, pady=10)

# Iniciar el loop de la ventana
root.mainloop()



if __name__ == '__main__':
    # Cambiar a 0.0.0.0 permite conexiones externas
    app.run(host='192.168.0.3', port=5000, debug=True)